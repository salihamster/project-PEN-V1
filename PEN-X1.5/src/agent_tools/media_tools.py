"""
WhatsApp Media Analysis Tools

Single tool for analyzing any WhatsApp media:
- Images → Gemini Vision
- PDFs → OCR / text extraction
- PPTX/DOCX → Document parsing
- Audio → (future: transcription)
"""

import os
import json
import base64
from pathlib import Path
from typing import Dict, Any, Optional

from ..parsers.media_manager import get_media_manager
from ..utils.logger import get_logger

logger = get_logger(__name__)


class MediaTools:
    """Tools for analyzing WhatsApp media files."""
    
    def __init__(self, gemini_api_key: Optional[str] = None):
        """
        Initialize MediaTools.
        
        Args:
            gemini_api_key: Gemini API key for vision analysis
        """
        self.gemini_api_key = gemini_api_key or os.getenv("GEMINI_API_KEY")
        self.media_manager = get_media_manager()
    
    def analyze_media(self, media_id: str, force_reprocess: bool = False) -> Dict[str, Any]:
        """
        Analyze a WhatsApp media file by its ID.
        
        Automatically determines file type and uses appropriate analyzer:
        - Images: Gemini Vision
        - PDFs: Text extraction
        - PPTX: Slide text extraction
        - DOCX: Document text extraction
        
        Results are cached to avoid re-processing.
        
        Args:
            media_id: Media ID (e.g., "IMG-20260103-WA0014")
            force_reprocess: If True, ignore cache and reprocess
            
        Returns:
            Analysis result with status and content
        """
        # Check if media exists
        media_info = self.media_manager.get_media_info(media_id)
        if not media_info:
            return {
                "status": "error",
                "error": "media_not_found",
                "message": f"Media ID '{media_id}' bulunamadi. Mevcut medyalari listelemek icin list_chat_media aracini kullanin."
            }
        
        # Check cache (unless force_reprocess)
        if not force_reprocess:
            cached = self.media_manager.get_cached_result(media_id)
            if cached:
                logger.info(f"Returning cached result for {media_id}")
                return {
                    "status": "success",
                    "media_id": media_id,
                    "media_type": media_info.get("type"),
                    "original_name": media_info.get("original_name"),
                    "chat_name": media_info.get("chat_name"),
                    "analysis": cached.get("result"),
                    "from_cache": True,
                    "processed_at": cached.get("processed_at")
                }
        
        # Get file path
        file_path = self.media_manager.get_media_path(media_id)
        if not file_path or not file_path.exists():
            return {
                "status": "error",
                "error": "file_not_found",
                "message": f"Medya dosyasi bulunamadi: {media_id}"
            }
        
        # Process based on type
        media_type = media_info.get("type", "unknown")
        
        try:
            if media_type == "image":
                result = self._analyze_image(file_path, media_info)
            elif media_type == "pdf":
                result = self._analyze_pdf(file_path, media_info)
            elif media_type == "pptx":
                result = self._analyze_pptx(file_path, media_info)
            elif media_type == "docx":
                result = self._analyze_docx(file_path, media_info)
            elif media_type == "audio":
                result = self._analyze_audio(file_path, media_info)
            elif media_type == "video":
                result = {
                    "status": "unsupported",
                    "message": f"Video dosyalari henuz desteklenmiyor: {media_id}"
                }
            else:
                result = {
                    "status": "unsupported",
                    "message": f"Bilinmeyen dosya tipi: {media_type}"
                }
            
            # Cache successful results
            if result.get("status") == "success" and "analysis" in result:
                self.media_manager.set_cached_result(
                    media_id,
                    result["analysis"],
                    model_used=result.get("model_used", "")
                )
            
            return result
            
        except Exception as e:
            logger.error(f"Error analyzing {media_id}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(type(e).__name__),
                "message": f"Analiz hatasi: {str(e)}"
            }
            
    def _analyze_audio(self, file_path: Path, media_info: Dict) -> Dict[str, Any]:
        """Analyze audio file using Gemini API."""
        if not self.gemini_api_key:
            return {
                "status": "error",
                "error": "no_api_key",
                "message": "Gemini API key bulunamadi"
            }
        
        import google.generativeai as genai
        genai.configure(api_key=self.gemini_api_key)
        
        # Read audio
        with open(file_path, 'rb') as f:
            audio_data = f.read()
        
        # Determine mime type
        ext = file_path.suffix.lower()
        mime_map = {
            '.mp3': 'audio/mp3',
            '.wav': 'audio/wav',
            '.ogg': 'audio/ogg',
            '.opus': 'audio/ogg', # Gemini supports ogg, opus is often in ogg container
            '.m4a': 'audio/mp4',
            '.aac': 'audio/aac'
        }
        mime_type = mime_map.get(ext, 'audio/mp3')
        
        # Send to Gemini
        model = genai.GenerativeModel('gemini-2.0-flash')
        
        prompt = """Bu ses dosyasını dinle ve şunları yap:
1. Konuşmanın tam transkriptini çıkar (dediklerini yaz).
2. Konuşmacının duygusunu veya tonunu kısaca belirt.
3. Önemli bir bilgi veya talimat varsa özetle.

Türkçe yaz."""
        
        response = model.generate_content([
            prompt,
            {'mime_type': mime_type, 'data': base64.b64encode(audio_data).decode()}
        ])
        
        return {
            "status": "success",
            "media_id": media_info.get("id"),
            "media_type": "audio",
            "original_name": media_info.get("original_name"),
            "chat_name": media_info.get("chat_name"),
            "sender": media_info.get("sender"),
            "timestamp": media_info.get("timestamp"),
            "analysis": response.text,
            "model_used": "gemini-2.0-flash (audio)",
            "from_cache": False
        }
    
    def _analyze_image(self, file_path: Path, media_info: Dict) -> Dict[str, Any]:
        """Analyze image using Gemini Vision."""
        if not self.gemini_api_key:
            return {
                "status": "error",
                "error": "no_api_key",
                "message": "Gemini API key bulunamadi"
            }
        
        import google.generativeai as genai
        genai.configure(api_key=self.gemini_api_key)
        
        # Read image
        with open(file_path, 'rb') as f:
            image_data = f.read()
        
        # Determine mime type
        ext = file_path.suffix.lower()
        mime_map = {
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.gif': 'image/gif',
            '.webp': 'image/webp'
        }
        mime_type = mime_map.get(ext, 'image/jpeg')
        
        # Send to Gemini Vision
        model = genai.GenerativeModel('gemini-2.0-flash')
        
        prompt = """Bu gorseli analiz et ve icerigini acikla:
1. Gorselde ne var? (metin, diyagram, fotograf, vs.)
2. Eger metin varsa, metni cikart
3. Onemli bilgileri ozetle

Turkce yaz, acik ve anlasilir ol."""
        
        response = model.generate_content([
            prompt,
            {'mime_type': mime_type, 'data': base64.b64encode(image_data).decode()}
        ])
        
        return {
            "status": "success",
            "media_id": media_info.get("id"),
            "media_type": "image",
            "original_name": media_info.get("original_name"),
            "chat_name": media_info.get("chat_name"),
            "sender": media_info.get("sender"),
            "timestamp": media_info.get("timestamp"),
            "analysis": response.text,
            "model_used": "gemini-2.0-flash",
            "from_cache": False
        }
    
    def _analyze_pdf(self, file_path: Path, media_info: Dict) -> Dict[str, Any]:
        """Analyze PDF file."""
        try:
            import fitz  # PyMuPDF
            
            doc = fitz.open(str(file_path))
            text_parts = []
            
            for page_num, page in enumerate(doc, 1):
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"[Sayfa {page_num}]\n{text.strip()}")
            
            doc.close()
            
            full_text = "\n\n".join(text_parts) if text_parts else "(PDF'den metin cikarildi)"
            
            return {
                "status": "success",
                "media_id": media_info.get("id"),
                "media_type": "pdf",
                "original_name": media_info.get("original_name"),
                "chat_name": media_info.get("chat_name"),
                "page_count": len(text_parts),
                "analysis": full_text[:10000],  # Limit to 10K chars
                "model_used": "pymupdf",
                "from_cache": False
            }
            
        except ImportError:
            return {
                "status": "error",
                "error": "library_missing",
                "message": "PyMuPDF kurulu degil. pip install pymupdf"
            }
    
    def _analyze_pptx(self, file_path: Path, media_info: Dict) -> Dict[str, Any]:
        """Analyze PowerPoint file."""
        try:
            from pptx import Presentation
            
            prs = Presentation(str(file_path))
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    slides_text.append(f"[Slayt {slide_num}]\n" + "\n".join(slide_content))
            
            full_text = "\n\n".join(slides_text) if slides_text else "(Slaytlardan metin cikarildi)"
            
            return {
                "status": "success",
                "media_id": media_info.get("id"),
                "media_type": "pptx",
                "original_name": media_info.get("original_name"),
                "chat_name": media_info.get("chat_name"),
                "slide_count": len(prs.slides),
                "analysis": full_text[:10000],
                "model_used": "python-pptx",
                "from_cache": False
            }
            
        except ImportError:
            return {
                "status": "error",
                "error": "library_missing",
                "message": "python-pptx kurulu degil. pip install python-pptx"
            }
    
    def _analyze_docx(self, file_path: Path, media_info: Dict) -> Dict[str, Any]:
        """Analyze Word document."""
        try:
            from docx import Document
            
            doc = Document(str(file_path))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            
            full_text = "\n\n".join(paragraphs) if paragraphs else "(Dokumandan metin cikarildi)"
            
            return {
                "status": "success",
                "media_id": media_info.get("id"),
                "media_type": "docx",
                "original_name": media_info.get("original_name"),
                "chat_name": media_info.get("chat_name"),
                "paragraph_count": len(paragraphs),
                "analysis": full_text[:10000],
                "model_used": "python-docx",
                "from_cache": False
            }
            
        except ImportError:
            return {
                "status": "error",
                "error": "library_missing",
                "message": "python-docx kurulu degil. pip install python-docx"
            }
    
    def list_chat_media(self, chat_name: str) -> Dict[str, Any]:
        """
        List all media files for a specific chat.
        
        Args:
            chat_name: WhatsApp chat name
            
        Returns:
            List of media with IDs and basic info
        """
        media_list = self.media_manager.list_media_by_chat(chat_name)
        
        if not media_list:
            return {
                "status": "success",
                "chat_name": chat_name,
                "total_media": 0,
                "media": [],
                "message": f"'{chat_name}' sohbetinde medya bulunamadi"
            }
        
        # Simplify for output
        simplified = []
        for m in media_list:
            simplified.append({
                "id": m.get("id"),
                "type": m.get("type"),
                "name": m.get("original_name"),
                "sender": m.get("sender"),
                "timestamp": m.get("timestamp"),
                "cached": m.get("id") in self.media_manager.cache
            })
        
        return {
            "status": "success",
            "chat_name": chat_name,
            "total_media": len(simplified),
            "media": simplified
        }
    
    def get_media_stats(self) -> Dict[str, Any]:
        """Get overall media statistics."""
        stats = self.media_manager.get_statistics()
        return {
            "status": "success",
            **stats
        }
